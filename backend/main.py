"""
Presentation Script Generator Backend

Analyzes PDF and PowerPoint presentations and generates speaker scripts
using a local LLM (Llama).
"""

import os
import re
import json
from typing import Optional, List, Dict, Any
from pathlib import Path
from io import BytesIO

from fastapi import FastAPI, File, UploadFile, HTTPException, Form
from fastapi.responses import JSONResponse, StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field

import pdfplumber
from pptx import Presentation
from llama_cpp import Llama
import uvicorn

# Initialize FastAPI app
app = FastAPI(title="Presentation Script Generator")

# CORS configuration
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Llama model cache
llm: Optional[Llama] = None

# Configuration
LLAMA_MODEL_PATH = os.getenv("LLAMA_MODEL_PATH", "./models/mistral-7b-instruct.Q4_K_M.gguf")
DEFAULT_MAX_TOKENS = 1000


class ScriptConfig(BaseModel):
    """Configuration for script generation."""
    duration_minutes: int = Field(default=5, ge=1, le=60, description="Target duration in minutes")
    tone: str = Field(default="professional", description="Tone of the script")
    focus_notes: Optional[str] = Field(default=None, description="Specific topics to focus on")
    include_notes: bool = Field(default=True, description="Whether to include speaker notes")
    output_format: str = Field(default="markdown", description="Output format")


def get_llm_model():
    """Load or return cached Llama model."""
    global llm
    if llm is None:
        if not os.path.exists(LLAMA_MODEL_PATH):
            raise HTTPException(
                status_code=500,
                detail=f"Llama model not found at {LLAMA_MODEL_PATH}. Please set LLAMA_MODEL_PATH or download a model."
            )
        llm = Llama(
            model_path=LLAMA_MODEL_PATH,
            n_ctx=4096,
            n_threads=4,
            n_gpu_layers=0  # CPU only
        )
    return llm


def extract_text_from_pdf(file_path: str) -> str:
    """Extract text content from a PDF file."""
    text_content = []
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                text_content.append(text)
    return "\n\n".join(text_content)


def extract_text_from_pptx(file_path: str) -> str:
    """Extract text content from a PowerPoint file."""
    prs = Presentation(file_path)
    text_content = []

    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            text_content.append("\n".join(slide_text))

    return "\n\n".join(text_content)


def analyze_presentation(content: str, config: ScriptConfig) -> Dict[str, Any]:
    """Analyze presentation content and extract key points."""
    prompt = f"""
    Analyze the following presentation content and extract:
    1. Main title/topic
    2. Key sections/outline
    3. Most important points
    4. Supporting details

    Content:
    {content[:2000]}

    Return as JSON with keys: title, sections, key_points, supporting_details.
    """

    try:
        llm = get_llm_model()
        response = llm(prompt, max_tokens=512, stop=["\n\n"])

        # Handle different response formats from llama-cpp-python
        if isinstance(response, dict):
            if "choices" in response:
                analysis_text = response["choices"][0]["text"] if isinstance(response["choices"], list) else response["choices"]
            else:
                analysis_text = str(response)
        else:
            analysis_text = str(response)

        # Try to parse as JSON, fall back to simple extraction if needed
        try:
            return json.loads(analysis_text)
        except json.JSONDecodeError:
            return {
                "title": "Presentation Analysis",
                "sections": ["Introduction", "Key Points", "Conclusion"],
                "key_points": ["Main point 1", "Main point 2"],
                "supporting_details": analysis_text[:500]
            }
    except Exception as e:
        return {
            "title": "Presentation",
            "sections": [],
            "key_points": [],
            "supporting_details": f"Error analyzing presentation: {str(e)}"[:500]
        }


def generate_script(content: str, config: ScriptConfig, analysis: Dict) -> str:
    """Generate a speaker script based on presentation content."""

    # Build context from analysis
    title = analysis.get("title", "Presentation")
    sections = analysis.get("sections", [])
    key_points = analysis.get("key_points", [])

    # Format focus notes
    focus_context = ""
    if config.focus_notes:
        focus_context = f"\n\nUSER-SPECIFIED FOCUS AREAS:\n{config.focus_notes}"

    # Build prompt
    prompt = f"""
You are an expert presentation coach. Generate a speaker script for a presentation titled "{title}".
Target duration: {config.duration_minutes} minutes
Tone: {config.tone}

Presentation Analysis:
- Sections: {', '.join(sections) if sections else 'Not available'}
- Key Points: {', '.join(key_points) if key_points else 'Not available'}

Content from slides:
{content[:3000]}

{focus_context}

Generate a comprehensive speaker script with:
1. An engaging opening that hooks the audience
2. Clear section transitions
3. Natural, conversational language appropriate for {config.tone} tone
4. {config.duration_minutes * 150} words (approx. 150 words per minute)
5. Speaker notes in parentheses where appropriate

Format as markdown with:
- # Title for each section
- Clear speaker dialogue
- (Notes) for delivery guidance

Script:
"""

    try:
        llm = get_llm_model()
        # Calculate tokens needed: ~150 words per minute * 4 tokens per word
        needed_tokens = config.duration_minutes * 150 * 4
        # Use a safe maximum between 8000 and the needed tokens
        max_tokens_calc = max(8000, needed_tokens)

        response = llm(
            prompt,
            max_tokens=max_tokens_calc,
            stop=[],
            stream=False
        )

        # Handle different response formats from llama-cpp-python
        if isinstance(response, dict):
            if "choices" in response and isinstance(response["choices"], list):
                return response["choices"][0]["text"]
            elif isinstance(response, dict):
                return str(response.get("choices", ""))
        return str(response)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error generating script: {str(e)}")


@app.post("/api/generate-script")
async def generate_script_endpoint(
    file: UploadFile = File(...),
    duration_minutes: int = Form(5),
    tone: str = Form("professional"),
    focus_notes: Optional[str] = Form(None),
    include_notes: bool = Form(True),
    output_format: str = Form("markdown")
):
    """Generate a presentation script from an uploaded file."""

    # Validate file type
    allowed_types = [".pdf", ".pptx", ".ppt"]
    file_extension = Path(file.filename).suffix.lower()

    if file_extension not in allowed_types:
        raise HTTPException(
            status_code=400,
            detail=f"Invalid file type. Please upload a PDF or PowerPoint file."
        )

    # Save uploaded file temporarily
    temp_path = f"temp_{file.filename}"
    try:
        with open(temp_path, "wb") as f:
            content = await file.read()
            f.write(content)

        # Extract text based on file type
        if file_extension == ".pdf":
            presentation_text = extract_text_from_pdf(temp_path)
        else:
            presentation_text = extract_text_from_pptx(temp_path)

        # Create config
        config = ScriptConfig(
            duration_minutes=duration_minutes,
            tone=tone,
            focus_notes=focus_notes,
            include_notes=include_notes,
            output_format=output_format
        )

        # Analyze presentation
        analysis = analyze_presentation(presentation_text, config)

        # Generate script
        script = generate_script(presentation_text, config, analysis)

        return {
            "success": True,
            "title": analysis.get("title", "Presentation"),
            "analysis": analysis,
            "script": script,
            "word_count": len(script.split()),
            "duration_minutes": duration_minutes
        }

    finally:
        # Clean up temp file
        if os.path.exists(temp_path):
            os.remove(temp_path)


@app.get("/api/health")
async def health_check():
    """Health check endpoint."""
    return {"status": "healthy"}


if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser()
    parser.add_argument("--host", default="0.0.0.0")
    parser.add_argument("--port", type=int, default=8000)
    args = parser.parse_args()

    uvicorn.run(app, host=args.host, port=args.port)
